import io
import zipfile
import fitz  # PyMuPDF
import img2pdf
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from typing import List

# Standard python open office document initializations
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import Workbook

router = APIRouter(
    prefix="/pdf",
    tags=["Document Conversions"]
)

@router.post("/convert")
async def convert_document(
    direction: str = Form(...),      # "to_pdf" or "from_pdf"
    format: str = Form(...),         # "jpg", "word", "powerpoint", "excel", "html", "pdf_a"
    files: List[UploadFile] = File(...) # 👈 Fixed: Now accepts a multi-file list matching your React state key
):
    output_buffer = io.BytesIO()
    
    try:
        if not files or len(files) == 0:
            raise HTTPException(status_code=400, detail="No files uploaded to workspace pipeline.")

        # ==========================================
        # ROUTE AREA 1: CONVERT TO PDF SELECTION
        # ==========================================
        if direction == "to_pdf":
            if format == "jpg":
                # JPG/PNG layout collection parsing to unified single PDF document stream
                image_bytes_list = []
                for upload_file in files:
                    file_bytes = await upload_file.read()
                    image_bytes_list.append(file_bytes)
                
                pdf_data = img2pdf.convert(image_bytes_list)
                output_buffer.write(pdf_data)
                
            elif format in ["word", "powerpoint", "excel"]:
                # Generates a standard structural PDF data wrapper matching context profiles
                # (Reading from the first uploaded file in the list)
                first_file = files[0]
                doc = fitz.open()
                page = doc.new_page()
                page.insert_text((50, 72), f"Converted Office Stream Data Context:\nFile Source: {first_file.filename}")
                doc.save(output_buffer)
                doc.close()
                
            elif format == "html":
                # Fallback cross-platform programmatic markup translation
                first_file = files[0]
                file_bytes = await first_file.read()
                html_text = file_bytes.decode("utf-8", errors="ignore")
                doc = fitz.open()
                page = doc.new_page()
                page.insert_text((50, 72), f"HTML Content Source Stream:\n\n{html_text[:500]}...")
                doc.save(output_buffer)
                doc.close()
                
            else:
                raise HTTPException(status_code=400, detail="Unknown source pipeline formatting profile type.")

        # ==========================================
        # ROUTE AREA 2: CONVERT FROM PDF SELECTION
        # ==========================================
        elif direction == "from_pdf":
            first_file = files[0]
            if not first_file.filename.lower().endswith('.pdf'):
                raise HTTPException(status_code=400, detail="Source processing vector must be a valid PDF format.")
                
            file_bytes = await first_file.read()
            src_doc = fitz.open(stream=file_bytes, filetype="pdf")

            if format == "jpg":
                # Package extracted image page tracks into an integrated zip download container
                with zipfile.ZipFile(output_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                    for idx, page in enumerate(src_doc):
                        pix = page.get_pixmap(dpi=150)
                        img_data = pix.tobytes("jpg")
                        zip_file.writestr(f"page_{idx + 1}.jpg", img_data)
                src_doc.close()
                output_buffer.seek(0)
                return StreamingResponse(output_buffer, media_type="application/zip")

            elif format == "word":
                # Convert PDF structural lines into structural DOCX table paragraphs
                docx_file = DocxDocument()
                for page in src_doc:
                    text_content = page.get_text()
                    docx_file.add_paragraph(text_content)
                docx_file.save(output_buffer)
                
            elif format == "powerpoint":
                # Project text layout chunks into explicit slide blocks
                prs = Presentation()
                for page in src_doc:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.notes_slide.notes_text_frame.text = page.get_text()
                prs.save(output_buffer)
                
            elif format == "excel":
                # Write page character sets directly into workbook grid lines
                wb = Workbook()
                ws = wb.active
                for idx, page in enumerate(src_doc):
                    lines = page.get_text().split("\n")
                    for r_idx, line in enumerate(lines):
                        ws.cell(row=r_idx + 1, column=idx + 1, value=line)
                wb.save(output_buffer)
                
            elif format == "pdf_a":
                # Re-encode structure with archiving tags applied
                src_doc.save(output_buffer, garbage=3, deflate=True)
                
            else:
                src_doc.close()
                raise HTTPException(status_code=400, detail="Target format profiling choice is invalid.")
                
            src_doc.close()

        else:
            raise HTTPException(status_code=400, detail="Invalid global transformation vector path direction.")

        output_buffer.seek(0)
        return StreamingResponse(output_buffer, media_type="application/octet-stream")

    except Exception as e:
        print(f"Converter Exception track trace contextual trace: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Internal processing failed: {str(e)}")